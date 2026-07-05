"""
Attachment text extraction helpers.

Parses PDF, DOCX, PPTX, XLSX, and image attachments into plain text
for downstream LLM extraction.
"""

from __future__ import annotations

import io
import logging
from typing import Literal

import pdfplumber
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image
import pytesseract

_TESSERACT_AVAILABLE = False
try:
    pytesseract.get_tesseract_version()
    _TESSERACT_AVAILABLE = True
except Exception:
    pass

logger = logging.getLogger(__name__)

AttachmentType = Literal["pdf", "docx", "pptx", "xlsx", "png", "jpg", "txt", "unknown"]


def _detect_type(filename: str, content_type: str | None = None) -> AttachmentType:
    """Detect attachment type from filename and optional MIME type."""
    fname_lower = filename.lower()
    if fname_lower.endswith(".pdf"):
        return "pdf"
    if fname_lower.endswith(".docx"):
        return "docx"
    if fname_lower.endswith(".pptx"):
        return "pptx"
    if fname_lower.endswith(".xlsx"):
        return "xlsx"
    if fname_lower.endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff")):
        ext = fname_lower.split(".")[-1]
        if ext in ("jpg", "jpeg"):
            return "jpg"
        return "png"  # catch-all for images
    if fname_lower.endswith(".txt"):
        return "txt"

    if content_type:
        ct = content_type.lower()
        if "pdf" in ct:
            return "pdf"
        if "wordprocessingml" in ct or "msword" in ct:
            return "docx"
        if "presentationml" in ct or "powerpoint" in ct:
            return "pptx"
        if "spreadsheetml" in ct or "excel" in ct:
            return "xlsx"
        if "image" in ct:
            if "jpeg" in ct or "jpg" in ct:
                return "jpg"
            return "png"
        if "text/plain" in ct:
            return "txt"

    return "unknown"


def _extract_pdf_text(file_bytes: bytes) -> str:
    """Extract text from a PDF using pdfplumber."""
    text_parts: list[str] = []
    try:
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
    except Exception as e:
        logger.error("PDF extraction failed: %s", e)
    return "\n".join(text_parts)


def _extract_docx_text(file_bytes: bytes) -> str:
    """Extract text from a DOCX file."""
    try:
        doc = DocxDocument(io.BytesIO(file_bytes))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract text from tables
        table_texts = []
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                table_texts.append(" | ".join(row_text))
        return "\n".join(paragraphs + table_texts)
    except Exception as e:
        logger.error("DOCX extraction failed: %s", e)
        return ""


def _extract_pptx_text(file_bytes: bytes) -> str:
    """Extract text from a PPTX presentation."""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_texts.append(shape.text)
            if slide_texts:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error("PPTX extraction failed: %s", e)
        return ""


def _extract_xlsx_text(file_bytes: bytes) -> str:
    """Extract text from an XLSX spreadsheet."""
    try:
        wb = load_workbook(io.BytesIO(file_bytes), data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_texts = [f"--- Sheet: {sheet_name} ---"]
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                if any(row_text):
                    sheet_texts.append(" | ".join(row_text))
            text_parts.append("\n".join(sheet_texts))
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error("XLSX extraction failed: %s", e)
        return ""


def _extract_image_text(file_bytes: bytes) -> str:
    """Extract text from an image using OCR (Tesseract)."""
    if not _TESSERACT_AVAILABLE:
        logger.warning("Tesseract OCR not available, skipping image text extraction")
        return ""
    try:
        image = Image.open(io.BytesIO(file_bytes))
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        logger.error("Image OCR failed: %s", e)
        return ""


def extract_text_from_attachment(
    file_bytes: bytes,
    filename: str,
    content_type: str | None = None,
) -> tuple[AttachmentType, str]:
    """
    Extract plain text from an attachment.

    Returns:
        (detected_type, extracted_text)
    """
    att_type = _detect_type(filename, content_type)

    if att_type == "pdf":
        text = _extract_pdf_text(file_bytes)
    elif att_type == "docx":
        text = _extract_docx_text(file_bytes)
    elif att_type == "pptx":
        text = _extract_pptx_text(file_bytes)
    elif att_type == "xlsx":
        text = _extract_xlsx_text(file_bytes)
    elif att_type in ("png", "jpg"):
        text = _extract_image_text(file_bytes)
    elif att_type == "txt":
        try:
            text = file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            text = file_bytes.decode("utf-8", errors="ignore")
    else:
        # Unknown type: try UTF-8 decode as fallback
        try:
            text = file_bytes.decode("utf-8", errors="ignore")
        except Exception:
            text = ""

    # Truncate to reasonable length for Kafka messages
    MAX_TEXT_LENGTH = 100_000  # ~100KB
    if len(text) > MAX_TEXT_LENGTH:
        text = text[:MAX_TEXT_LENGTH] + "\n...[truncated]"

    logger.info(
        "Extracted %d chars from %s (type: %s)", len(text), filename, att_type
    )
    return att_type, text
